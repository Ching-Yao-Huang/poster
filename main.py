
import os, sys, math
import pandas as pd

from pptx import Presentation
from pptx.util import Mm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


eng_font_style = ["微軟正黑體", "Brush Script MT"]
ch_font_style = ["微軟正黑體", "標楷體"]
font_default_size = {
    "微軟正黑體": Mm(34), \
    "Brush Script MT": Mm(50), \
    "標楷體": Mm(38)
}

def add_text(slide, left, top, width, height, text, **kwargs):
    tbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = tbox.text_frame
    text_frame.word_wrap = True

    p = text_frame.add_paragraph()
    p.text = text

    assert "size" in kwargs
    p.font.size = kwargs["size"]
    p.font.bold = True

    assert "style" in kwargs
    p.font.name = kwargs["style"]

    if "center" in kwargs and kwargs["center"]:
        p.alignment = PP_ALIGN.CENTER
    p.font.color.rgb = RGBColor(0, 102, 204)

def style_config(text, style):

    ch = False
    eng = False
    for c in text:
        if c.isdigit():
            continue
        if c in "abcdefghijklmnopqrstuvwxyzABCDEFGHIJKLMNOPQRSTUVWXYZ":
            eng = True
        else:
            ch = True
    
    if eng and ch:
        return "微軟正黑體", font_default_size["微軟正黑體"]
    else:
        if eng:
            return eng_font_style[style], font_default_size[eng_font_style[style]]
        else:
            return ch_font_style[style], font_default_size[ch_font_style[style]]

def layout1(W, H, data, slide, r1, r2, bbox_height):
    h = 0
    for i, s_data in enumerate(data[2: ]):
        if i == 0:
            # print(s_data)
            style, size = style_config("con", s_data[0])
            print(style, size)
            add_text(slide, \
                W / 8, H / 10, \
                W, bbox_height, \
                "Congratulations", size = size, style = style)
            h = H / 5
        if i == 1:
            style, size = style_config(s_data[0] + s_data[1] + "榮獲", s_data[2])
            # print("section 1", style, size)
            # ----------------------------------------------
            N = math.floor((4 * W / 5) / size) + 1
            N_Ls = math.ceil((len(s_data[0]) + 2) / N)
            # print(N, N_Ls)
            add_text(slide, \
                W / 5, h, \
                4 * W / 5, N_Ls * bbox_height, \
                "賀！" + s_data[0], size = size, style = style)
            h += N_Ls * bbox_height

            # ----------------------------------------------
            add_text(slide, \
                2 * W / 5, h, \
                3 * W / 5 , bbox_height, \
                "榮獲", size = size, style = style)
            h += bbox_height

            # ----------------------------------------------
            size *= r1
            N = (W - W / 25) / size
            N_Ls = math.ceil(len(s_data[1]) / N)
            add_text(slide, \
                W / 25, h, \
                W - W / 25, N_Ls * bbox_height * r1, \
                s_data[1], size = size, style = style)
            h += N_Ls * bbox_height * r1
        
        if i == 2 or i == 3 or i == 4:
            style, size = style_config(s_data[0] + "榮獲", s_data[1])
            print(style, size)
            size *= r1

            if i == 2:
                pre = "獲獎隊伍："
            elif i == 3:
                pre = "獲獎學生："
            else:
                pre = "指導教授："

            N = (W - W / 25) / size
            N_Ls = math.ceil((len(s_data[0]) + 5) / N)
                
            add_text(slide, \
                W / 25, h, \
                W - W / 25, N_Ls * bbox_height * r1, \
                pre + s_data[0], size = size, style = style)
            h += N_Ls * bbox_height * r1

        if i == 5:
            h += H / 20
            slide.shapes.add_picture(s_data[0], W / 4, h, width = W / 2, )

            min_h = H + 1
            for shape in slide.shapes:
                if shape.shape_type == 13:
                    height = shape.height
                    min_h = min(min_h, height)
            h += min_h

        if i == 6:
            style, size = style_config("榮獲", s_data[0])
            print(style, size)
            size *= r2
            add_text(slide, \
                W / 2, h, \
                W / 2, bbox_height * r2, \
                "資工系全體師生祝賀", size = size, style = style)

def layout2(W, H, data, slide, r1, r2, bbox_height):
    h = 0
    img_h = None
    for i, s_data in enumerate(data[2: ]):
        if i == 0:
            # print(s_data)
            style, size = style_config("con", s_data[0])
            print(style, size)
            add_text(slide, \
                0, 7 * H / 24, \
                W, bbox_height, \
                "Congratulations", size = size, style = style, center = True)
            h = H / 2
            img_h = h
        if i == 1:
            style, size = style_config(s_data[0] + s_data[1] + "榮獲", s_data[2])
            # print("section 1", style, size)
            # ----------------------------------------------
            N = math.floor((W / 2 - W / 25) / size) + 1
            N_Ls = math.ceil((len(s_data[0]) + 2) / N)
            # print(N, N_Ls)
            add_text(slide, \
                W / 25, h, \
                W / 2 - W / 25, N_Ls * bbox_height, \
                "賀！" + s_data[0], size = size, style = style)
            h += N_Ls * bbox_height

            # ----------------------------------------------
            add_text(slide, \
                W / 25, h, \
                W / 2 - W / 25 , bbox_height, \
                "榮獲", size = size, style = style)
            h += bbox_height

            # ----------------------------------------------
            N = math.floor((W / 2 - W / 25) / size) + 1
            N_Ls = math.ceil(len(s_data[1]) / N)
            add_text(slide, \
                W / 25, h, \
                W / 2 - W / 25, N_Ls * bbox_height, \
                s_data[1], size = size, style = style)
        
        if i == 2:
            slide.shapes.add_picture(s_data[0], W / 2, img_h, width = W / 2, )

        if i == 3:
            style, size = style_config("榮獲", s_data[0])
            print(style, size)
            size *= r2
            add_text(slide, \
                0, H - bbox_height * r2 - H / 100, \
                W, bbox_height * r2, \
                "資工系全體師生祝賀", size = size, style = style, center = True)

def main(W, H, datas):
    ppt = Presentation()

    ppt.slide_width = W
    ppt.slide_height = H

    r1 = 0.6
    r2 = 0.5
    bbox_height = H / 12
    for data in datas:

        slide_layout = ppt.slide_layouts[5]
        slide = ppt.slides.add_slide(slide_layout)
        for shape in slide.shapes:
            if not shape.has_text_frame:  # This ensures you're only removing text-related placeholders
                continue
            sp = shape.element
            sp.getparent().remove(sp)

        slide.shapes.add_picture(os.path.join("./bg", data[0]), 0, 0, width = W, height = H)
        layout_style = data[1]

        print(data)
        if layout_style == 1:
            layout1(W, H, data, slide, r1, r2, bbox_height)
        if layout_style == 2:
            layout2(W, H, data, slide, r1, r2, bbox_height)
        
    ppt.save('example.pptx')

def read_input(filename):
    df1 = pd.read_excel(filename, header = None)
    
    values = df1.values
    datas = []
    for row in values[2: ]:

        bg = row[0]
        layout = int(row[1])
        if layout == 1:
            s0_style = int(row[2])

            s1_text1 = row[3]
            s1_text2 = row[4]
            s1_style = int(row[5])

            s2_text = row[6]
            s2_style = int(row[7])

            s3_text = row[8]
            s3_style = int(row[9])

            s4_text = row[10]
            s4_style = int(row[11])
            
            img_path = os.path.join("./fig", row[12])
            s5_style = row[13]

            datas.append([
                bg, layout, [s0_style], \
                    [s1_text1, s1_text2, s1_style], \
                    [s2_text, s2_style], \
                    [s3_text, s3_style], \
                    [s4_text, s4_style], \
                    [img_path], \
                    [s5_style]
            ])
        else:
            s0_style = int(row[2])

            s1_text1 = row[3]
            s1_text2 = row[4]
            s1_style = int(row[5])

            img_path = os.path.join("./fig", row[12])
            s5_style = row[13]
            datas.append([
                bg, layout, [s0_style], \
                    [s1_text1, s1_text2, s1_style], \
                    [img_path], \
                    [s5_style]
            ])

    return datas

if __name__ == "__main__":

    datas = read_input(sys.argv[1])
    main(Mm(420), Mm(594), datas)